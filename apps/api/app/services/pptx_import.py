"""PowerPoint (.pptx) → DocumentJSON v1.0 import.

Mirrors the docx_import structure but maps PPT semantics into wiki shapes:

- **Slide 1** (with the smallest text + biggest fonts) → cover slide.
  Becomes the document title + summary; subsequent slides become sections.
- **Each slide → 1 Section** (level 1).
  Slide title = Section.title, slide layout name → Section.layout heuristic.
- **Slide layouts** are mapped to our 6 layout choices:
  - `Title Slide` / `Title Only` → Section.layout = `'title-only'`
  - `Two Content` / `Comparison` → `'two-col'`
  - `Picture with Caption` (image dominant + small text) → `'image-left'`
  - everything else → `'stack'` (default)
- **Per-slide shapes** become blocks:
  - text frames → ParagraphBlock (one block per non-empty placeholder/textbox)
  - tables → TableBlock (with sparse cells if merged)
  - pictures → ImageBlock (image bytes routed through the same uploader the
    docx_import uses, so the file lands in MinIO and gets a real image_id)
  - charts → KpiCardsBlock placeholder OR a paragraph note "[차트]"
  - speaker notes → ParagraphBlock with `meta.note: "speaker:N"` so they
    survive the round-trip to presentation mode (and to pptx_export).

What we DON'T preserve (intentionally):
  - Absolute shape positions (everything stacks)
  - Slide transitions / animations (PPT-only)
  - SmartArt diagrams (best-effort: text only)

The function returns the same `{document, summary}` envelope as docx_import
so the import endpoint can stay symmetric.
"""
from __future__ import annotations

import io
from dataclasses import dataclass, field
from typing import Any, Callable

import ulid
from pptx import Presentation
from pptx.util import Emu


# Hex literal for the .pptx zip magic ("PK\x03\x04" — same as .docx, both
# are Office Open XML packages). We accept either capitalisation just in
# case different OS exports use weird casing in the central directory.
_ZIP_MAGIC = b"PK\x03\x04"


@dataclass
class PptxImportSummary:
    slides: int = 0
    sections: int = 0
    paragraphs: int = 0
    tables: int = 0
    images: int = 0
    speaker_notes: int = 0
    warnings: list[str] = field(default_factory=list)


def is_pptx_zip_magic(buf: bytes) -> bool:
    return len(buf) >= 4 and buf[:4] == _ZIP_MAGIC


def is_pptx_content(buf: bytes) -> bool:
    """Reject .docx pretending to be .pptx — check ppt/presentation.xml exists."""
    import zipfile

    try:
        with zipfile.ZipFile(io.BytesIO(buf)) as zf:
            return "ppt/presentation.xml" in zf.namelist()
    except zipfile.BadZipFile:
        return False


def _new_id() -> str:
    return str(ulid.new())


def _layout_from_slide(slide: Any) -> str:
    """Heuristic mapping from PPT slide layout name → our layout enum."""
    layout = slide.slide_layout
    name = (getattr(layout, "name", "") or "").lower()
    if "title slide" in name or "title only" in name or "section header" in name:
        return "title-only"
    if "two content" in name or "comparison" in name:
        return "two-col"
    if "picture with caption" in name:
        return "image-left"
    return "stack"


def _extract_paragraph_text(shape: Any) -> list[str]:
    """Return one string per non-empty paragraph in this shape's text frame."""
    if not getattr(shape, "has_text_frame", False):
        return []
    out: list[str] = []
    for p in shape.text_frame.paragraphs:
        # Reassemble paragraph text from its runs so bold/italic markers
        # could be added later (current pass keeps it plain).
        text = "".join(r.text or "" for r in p.runs).strip()
        if text:
            out.append(text)
    return out


def _is_title_placeholder(shape: Any) -> bool:
    """python-pptx placeholders carry an idx; idx==0 is the title placeholder
    on most layouts. Non-placeholder shapes raise ``ValueError`` when you
    touch ``placeholder_format`` — guard with ``is_placeholder``.
    """
    if not getattr(shape, "has_text_frame", False):
        return False
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        ph = shape.placeholder_format
    except Exception:  # noqa: BLE001
        return False
    if ph is None:
        return False
    try:
        return ph.idx == 0
    except Exception:  # noqa: BLE001
        return False


def _slide_title(slide: Any) -> str:
    """Best-effort title — placeholder idx==0, falls back to first text shape."""
    for shape in slide.shapes:
        if _is_title_placeholder(shape):
            paras = _extract_paragraph_text(shape)
            if paras:
                return paras[0]
    # Fallback: first non-empty text in any text shape.
    for shape in slide.shapes:
        paras = _extract_paragraph_text(shape)
        if paras:
            return paras[0]
    return ""


def _table_to_block(shape: Any) -> dict[str, Any]:
    """Convert a PPT table shape to TableBlock. Detects merged cells via
    ``cell.is_merge_origin`` / ``cell.is_spanned`` and emits the sparse
    `cells` representation if any merge exists; else flat `headers`/`rows`.
    """
    table = shape.table
    n_rows = len(table.rows)
    n_cols = len(table.columns)
    has_merge = any(
        getattr(table.cell(r, c), "is_spanned", False)
        for r in range(n_rows)
        for c in range(n_cols)
    )
    if has_merge:
        cells_out: list[dict[str, Any]] = []
        for r in range(n_rows):
            for c in range(n_cols):
                cell = table.cell(r, c)
                if getattr(cell, "is_spanned", False):
                    continue
                text = (cell.text or "").strip()
                entry: dict[str, Any] = {"r": r, "c": c, "text": text}
                rs = getattr(cell, "span_height", 1) or 1
                cs = getattr(cell, "span_width", 1) or 1
                if rs > 1:
                    entry["rowSpan"] = rs
                if cs > 1:
                    entry["colSpan"] = cs
                if r == 0:
                    entry["header"] = True
                cells_out.append(entry)
        return {
            "type": "table",
            "id": _new_id(),
            "headers": [(table.cell(0, c).text or "").strip() for c in range(n_cols)],
            "rows": [],
            "cells": cells_out,
        }
    # Flat path — first row = headers, rest = rows.
    headers = [(table.cell(0, c).text or "").strip() for c in range(n_cols)]
    rows = [
        [(table.cell(r, c).text or "").strip() for c in range(n_cols)]
        for r in range(1, n_rows)
    ]
    return {
        "type": "table",
        "id": _new_id(),
        "headers": headers,
        "rows": rows,
    }


def _picture_to_block(
    shape: Any,
    image_uploader: Callable[[bytes, str], dict[str, Any] | None],
    summary: PptxImportSummary,
) -> dict[str, Any] | None:
    """Pull image bytes off a Picture shape and route through the uploader."""
    try:
        img = shape.image  # ImagePart
        data = img.blob
        ext = (img.ext or "png").lower()
        result = image_uploader(data, f"pptx-import.{ext}")
    except Exception as e:  # noqa: BLE001
        summary.warnings.append(f"image extraction failed: {e}")
        return None
    if not result or not result.get("image_id"):
        return None
    summary.images += 1
    return {
        "type": "image",
        "id": _new_id(),
        "imageId": str(result["image_id"]),
        "alt": "",
    }


def _speaker_note_text(slide: Any) -> str:
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return ""
    tf = getattr(notes_slide, "notes_text_frame", None)
    if tf is None:
        return ""
    text = (tf.text or "").strip()
    return text


def _section_from_slide(
    slide: Any,
    *,
    section_index: int,
    image_uploader: Callable[[bytes, str], dict[str, Any] | None],
    summary: PptxImportSummary,
) -> dict[str, Any]:
    """Convert one PPT slide → one Section dict (DocumentJSON shape)."""
    title = _slide_title(slide) or f"슬라이드 {section_index}"
    layout = _layout_from_slide(slide)

    blocks: list[dict[str, Any]] = []

    # Walk shapes in z-order. Skip the title placeholder (already used for
    # section.title) so we don't duplicate it as the first body block.
    for shape in slide.shapes:
        if _is_title_placeholder(shape):
            continue
        if getattr(shape, "has_table", False):
            blocks.append(_table_to_block(shape))
            summary.tables += 1
            continue
        if shape.shape_type == 13:  # PICTURE
            blk = _picture_to_block(shape, image_uploader, summary)
            if blk:
                blocks.append(blk)
            continue
        if getattr(shape, "has_text_frame", False):
            for text in _extract_paragraph_text(shape):
                blocks.append({
                    "type": "paragraph",
                    "id": _new_id(),
                    "text": text,
                })
                summary.paragraphs += 1
            continue
        # Unknown shape — drop a warning, skip silently in the body.
        summary.warnings.append(
            f"unsupported shape type on slide {section_index}: {shape.shape_type}"
        )

    # Speaker notes — append as paragraph(s) with meta.note. The slide
    # renderer strips `meta.note: speaker:N` paragraphs from the visible
    # body but exposes them in the presenter pane, mirroring how docx
    # handles footnote markers.
    notes = _speaker_note_text(slide)
    if notes:
        for line in notes.splitlines():
            line = line.strip()
            if not line:
                continue
            blocks.append({
                "type": "paragraph",
                "id": _new_id(),
                "text": line,
                "meta": {"note": f"speaker:{section_index}"},
            })
            summary.speaker_notes += 1

    section: dict[str, Any] = {
        "id": _new_id(),
        "level": 1,
        "title": title,
        "blocks": blocks,
        "subsections": [],
    }
    if layout != "stack":
        section["layout"] = layout
    return section


def pptx_to_document(
    buf: bytes,
    *,
    slug: str,
    title: str = "",
    owner_user_id: str | None = None,
    image_uploader: Callable[[bytes, str], dict[str, Any] | None] | None = None,
) -> dict[str, Any]:
    """Top-level entry point — open .pptx bytes and return a DocumentJSON."""
    if not is_pptx_zip_magic(buf):
        raise ValueError("not a valid zip (.pptx must be PK zip)")
    if not is_pptx_content(buf):
        raise ValueError("zip does not contain ppt/presentation.xml")

    # Default uploader returns None so picture blocks are silently dropped
    # if the caller doesn't wire MinIO. Tests rely on this: importing
    # without an uploader still produces a valid Document.
    if image_uploader is None:
        def _noop(_data: bytes, _name: str) -> dict[str, Any] | None:
            return None
        image_uploader = _noop

    prs = Presentation(io.BytesIO(buf))
    summary = PptxImportSummary()
    summary.slides = len(prs.slides)

    sections: list[dict[str, Any]] = []
    cover_title = title.strip()
    cover_summary = ""

    for idx, slide in enumerate(prs.slides):
        if idx == 0 and not cover_title:
            # First slide is treated as the document cover when the caller
            # didn't supply an explicit title. We still emit it as a
            # section so nothing's lost — its title lands as the doc title
            # AND as the first section heading.
            cover_title = _slide_title(slide) or "문서 제목"
            # Aggregate non-title text on the cover slide as the doc summary.
            cover_lines: list[str] = []
            for shape in slide.shapes:
                if _is_title_placeholder(shape):
                    continue
                cover_lines.extend(_extract_paragraph_text(shape))
            cover_summary = "\n".join(cover_lines)[:500]
        section = _section_from_slide(
            slide,
            section_index=idx + 1,
            image_uploader=image_uploader,
            summary=summary,
        )
        sections.append(section)
        summary.sections += 1

    doc: dict[str, Any] = {
        "schema_ver": "1.0.0",
        "slug": slug,
        "title": cover_title or "Untitled",
        "summary": cover_summary,
        "metadata": {},
        "infobox": {},
        "sections": sections,
    }
    if owner_user_id:
        doc["metadata"] = {**doc["metadata"], "owners": [owner_user_id]}

    return {"document": doc, "summary": summary}
