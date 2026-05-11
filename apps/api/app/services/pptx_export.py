"""DocumentJSON v1.0 → PowerPoint (.pptx) 출력.

Why python-pptx
---------------
- pure-python (no native deps) → apptainer 환경에서 추가 설정 없이 동작.
- ``Table`` / ``add_chart`` / ``add_picture`` 등 핵심 도형을 지원.
- WeasyPrint 와 달리 시스템 라이브러리 (cairo/pango) 가 필요 없다.

설계 요약
---------
- 슬라이드 1장 = level-1 섹션 1개.
- level-2 sub-section 은 본문 합계가 ≤ 500자면 inline bullets,
  초과하면 별도 sub-slide 로 분리 (``_SUBSECTION_INLINE_THRESHOLD``).
- block 별 핸들러는 :data:`_BLOCK_HANDLERS` 에 등록. 각 핸들러는
  현재 슬라이드의 본문 frame 에 paragraph/run 을 채워넣는다.
- 외부 의존성 (이미지/MinIO) 은 ``image_resolver`` callback 으로 주입.
  resolver 는 image_id → ``{bytes, mime}`` 를 돌려준다 (선택). 미지정/실패 시
  플레이스홀더 텍스트로 대체.
- ``meta.note`` 가 ``"speaker:"`` 로 시작하면 해당 슬라이드의 노트 페인에
  speaker note 로 누적 (I2 agent 의 presentation polish 와 약속한 prefix).
- ``meta.note == "page-break-before"`` 인 paragraph block 은 새 슬라이드를
  강제로 시작 (markdown / html 익스포트와 동일한 컨벤션).

부수효과 없음 — 입력 dict 변경하지 않으며, 결과는 BytesIO 에 쓴 .pptx 바이트.
"""
from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import Any, Callable

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.services.variables import walk_doc_substitute


# ── Tunables ─────────────────────────────────────────────────────────

# level-2 subsection 본문이 이 길이 이하면 inline bullets, 초과하면 sub-slide.
_SUBSECTION_INLINE_THRESHOLD = 500

# 슬라이드 좌표 (16:9, 13.333 × 7.5 in).
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)
_BODY_LEFT = Inches(0.5)
_BODY_TOP = Inches(1.4)
_BODY_WIDTH = Inches(12.333)
_BODY_HEIGHT = Inches(5.6)

_TITLE_TOP = Inches(0.4)
_TITLE_HEIGHT = Inches(0.9)

# Markdown-lite inline 토큰 — 굵게 / 기울임 / 취소선만 처리한다.
_INLINE_RE = re.compile(
    r"(\*\*([^*]+)\*\*|\*([^*]+)\*|~~([^~]+)~~)"
)

ImageResolver = Callable[[str], dict[str, Any] | None]


# ── Public entry ─────────────────────────────────────────────────────


@dataclass
class PptxOptions:
    """렌더 옵션.

    - ``image_resolver``: image_id → ``{bytes: bytes, mime: str}`` 또는 None.
    """

    image_resolver: ImageResolver | None = None


def render_pptx(
    doc: dict[str, Any],
    *,
    options: PptxOptions | None = None,
    requester_role: str | None = None,
) -> bytes:
    """DocumentJSON dict 를 .pptx 바이너리로 렌더한다.

    Args:
        doc: DocumentJSON v1.0 dict.
        options: 렌더 옵션. None 이면 기본값.
        requester_role: 호출자 role. 지정 시 admin 미만은 meta.permission 이
            높은 블록을 redact 한 결과를 렌더한다. None 이면 scrub 미적용.

    Returns:
        .pptx zip 바이너리 (PK\\x03\\x04 매직).
    """
    opts = options or PptxOptions()
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    ctx = _Ctx(prs=prs, opts=opts)

    if requester_role is not None:
        from .document_service import scrub_for_response

        doc = scrub_for_response(doc, role=requester_role)
    # Substitute `{{var}}` tokens up front so every slide emitter sees fully
    # resolved text. Code blocks are intentionally skipped inside the helper.
    doc = walk_doc_substitute(doc, doc.get("variables"))

    _render_title_slide(ctx, doc)

    for section in doc.get("sections") or []:
        _render_section(ctx, section)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Internal context ────────────────────────────────────────────────


@dataclass
class _Ctx:
    prs: Any
    opts: PptxOptions


# ── Title slide ──────────────────────────────────────────────────────


def _render_title_slide(ctx: _Ctx, doc: dict[str, Any]) -> None:
    """첫 슬라이드 — 문서 제목/요약/소속 라인."""
    slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])  # blank layout
    title = _str(doc.get("title")) or "Untitled"
    summary = _str(doc.get("summary"))
    metadata = doc.get("metadata") or {}

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True

    if summary:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(1.0))
        sf = sub_box.text_frame
        sf.word_wrap = True
        sp = sf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        run = sp.add_run()
        run.text = summary
        run.font.size = Pt(18)
        run.font.italic = True

    # division/team/group/part — slide footer line.
    parts = [
        _str(metadata.get(k))
        for k in ("division", "team", "group", "part")
        if _str(metadata.get(k))
    ]
    if parts:
        meta_box = slide.shapes.add_textbox(Inches(1), Inches(5.4), Inches(11.3), Inches(0.6))
        mf = meta_box.text_frame
        mp = mf.paragraphs[0]
        mp.alignment = PP_ALIGN.CENTER
        run = mp.add_run()
        run.text = " / ".join(parts)
        run.font.size = Pt(14)


# ── Section rendering ────────────────────────────────────────────────


def _render_section(ctx: _Ctx, section: dict[str, Any]) -> None:
    """level-1 섹션 = 슬라이드 1장 (+ optional sub-slides).

    level-2 subsection 은 분량에 따라 inline bullets 또는 별도 sub-slide.
    level-3 은 항상 inline bullets (depth 2) 로 합쳐진다 — sub-slide 분리는
    level-2 까지만.
    """
    number = _str(section.get("number"))
    title = _str(section.get("title"))
    heading = f"{number} {title}".strip()
    blocks = section.get("blocks") or []
    subsections = section.get("subsections") or []
    layout = _str(section.get("layout")) or "stack"

    slide = _new_content_slide(ctx, heading)

    # `title-only` is the cover-style layout — schema convention says the
    # body should be hidden. Skip rendering blocks but still emit the
    # slide so navigation order is preserved.
    if layout == "title-only":
        # Speaker notes (if any) still go through so the presenter view
        # picks them up; nothing else.
        for block in blocks:
            meta = block.get("meta") or {}
            note = _str(meta.get("note") if isinstance(meta, dict) else "")
            if note.startswith("speaker:"):
                speaker_text = note[len("speaker:") :].strip() or _str(block.get("text"))
                if speaker_text:
                    _append_speaker_note(slide, speaker_text)
        return

    body_frame = _body_text_frame(slide)

    # Render top-level blocks of this section.
    has_content = False
    for block in blocks:
        if _render_block_into(slide, body_frame, block, ctx, depth=0):
            has_content = True

    # Decide subsection placement.
    inline_subs: list[dict[str, Any]] = []
    spawn_subs: list[dict[str, Any]] = []
    for sub in subsections:
        sub_chars = _approx_char_count(sub)
        if sub_chars <= _SUBSECTION_INLINE_THRESHOLD:
            inline_subs.append(sub)
        else:
            spawn_subs.append(sub)

    for sub in inline_subs:
        sub_title = _str(sub.get("title"))
        sub_num = _str(sub.get("number"))
        heading_line = f"{sub_num} {sub_title}".strip()
        # Inline subsection heading as a bold paragraph.
        p = body_frame.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = heading_line
        run.font.bold = True
        run.font.size = Pt(20)
        for block in sub.get("blocks") or []:
            if _render_block_into(slide, body_frame, block, ctx, depth=1):
                has_content = True
        # level-3 subsections inside an inline level-2 collapse into depth 2.
        for deep in sub.get("subsections") or []:
            deep_title = _str(deep.get("title"))
            deep_num = _str(deep.get("number"))
            dp = body_frame.add_paragraph()
            dp.level = 1
            run = dp.add_run()
            run.text = f"{deep_num} {deep_title}".strip()
            run.font.bold = True
            run.font.size = Pt(16)
            for block in deep.get("blocks") or []:
                _render_block_into(slide, body_frame, block, ctx, depth=2)
        has_content = True

    if not has_content:
        # Empty body – placeholder hint so slide isn't blank.
        p = body_frame.paragraphs[0]
        run = p.add_run()
        run.text = "(내용 없음)"
        run.font.italic = True
        run.font.size = Pt(14)

    for sub in spawn_subs:
        _render_subsection_slide(ctx, sub)


def _render_subsection_slide(ctx: _Ctx, sub: dict[str, Any]) -> None:
    """level-2 subsection 이 길 때 별도 슬라이드로 분리."""
    number = _str(sub.get("number"))
    title = _str(sub.get("title"))
    heading = f"{number} {title}".strip()
    slide = _new_content_slide(ctx, heading)
    frame = _body_text_frame(slide)
    has_content = False
    for block in sub.get("blocks") or []:
        if _render_block_into(slide, frame, block, ctx, depth=0):
            has_content = True
    for deep in sub.get("subsections") or []:
        d_title = _str(deep.get("title"))
        d_num = _str(deep.get("number"))
        p = frame.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = f"{d_num} {d_title}".strip()
        run.font.bold = True
        run.font.size = Pt(18)
        for block in deep.get("blocks") or []:
            if _render_block_into(slide, frame, block, ctx, depth=1):
                has_content = True
    if not has_content:
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = "(내용 없음)"
        run.font.italic = True


# ── Slide / frame factories ──────────────────────────────────────────


def _new_content_slide(ctx: _Ctx, heading: str) -> Any:
    slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(_BODY_LEFT, _TITLE_TOP, _BODY_WIDTH, _TITLE_HEIGHT)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = heading
    run.font.size = Pt(28)
    run.font.bold = True
    return slide


def _body_text_frame(slide: Any) -> Any:
    body_box = slide.shapes.add_textbox(_BODY_LEFT, _BODY_TOP, _BODY_WIDTH, _BODY_HEIGHT)
    tf = body_box.text_frame
    tf.word_wrap = True
    # Default first paragraph stays empty until first block writes into it.
    return tf


# ── Block dispatch ───────────────────────────────────────────────────


def _render_block_into(
    slide: Any,
    frame: Any,
    block: dict[str, Any],
    ctx: _Ctx,
    *,
    depth: int = 0,
) -> bool:
    """Render a single block. Returns ``True`` if anything was emitted.

    The handler receives the slide (for shape additions like tables/images)
    and the body text frame (for paragraph appends). Some block types — chart,
    table, image — append shapes to the slide directly and return True.
    """
    btype = _str(block.get("type"))
    meta = block.get("meta") or {}

    # Honor the per-block audience opt-out — `wiki-only` blocks never
    # appear in slide exports. Mirrors the FE `filterForAudience` used by
    # the in-browser Presentation page so wiki/slide outputs stay aligned.
    if isinstance(meta, dict):
        audience = _str(meta.get("audience"))
        if audience == "wiki-only":
            return False

    # speaker note — append to slide notes pane regardless of type.
    if isinstance(meta, dict):
        note = _str(meta.get("note"))
        if note.startswith("speaker:"):
            speaker_text = note[len("speaker:") :].strip()
            if speaker_text:
                _append_speaker_note(slide, speaker_text)
        # page-break-before paragraph collapses to a no-op (the next block
        # would already start fresh on the next slide if section-level
        # boundaries are honoured). Inside a slide we just skip empty
        # carriers.
        if note == "page-break-before" and btype == "paragraph" and not _str(
            block.get("text")
        ):
            return False

    handler = _BLOCK_HANDLERS.get(btype)
    if handler is None:
        return _b_unknown(slide, frame, block, ctx, depth)
    try:
        return handler(slide, frame, block, ctx, depth)
    except Exception as e:  # pragma: no cover — defensive
        p = frame.add_paragraph()
        p.level = depth
        run = p.add_run()
        run.text = f"[블록 렌더 실패: {btype} — {e}]"
        run.font.size = Pt(12)
        run.font.italic = True
        return True


# ── Per-block handlers ───────────────────────────────────────────────


def _b_paragraph(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    text = _str(block.get("text"))
    if not text:
        return False
    p = _next_paragraph(frame)
    p.level = depth
    _emit_inline_runs(p, text, base_size=Pt(14))
    return True


def _b_heading_4(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    title = _str(block.get("title"))
    if not title:
        return False
    p = _next_paragraph(frame)
    p.level = depth
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(18)
    return True


def _b_list(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    style = _str(block.get("style")) or "bullet"
    items = block.get("items") or []
    if not items:
        return False
    for idx, item in enumerate(items, start=1):
        if isinstance(item, dict):
            text = _str(item.get("text"))
            item_depth = int(item.get("depth") or 0)
        else:
            text = _str(item)
            item_depth = 0
        if not text:
            continue
        prefix = ""
        if style == "number":
            prefix = f"{idx}. "
        elif style == "check":
            prefix = "☐ "
        # else bullet — frame's default bullet renders as • in PowerPoint.
        p = _next_paragraph(frame)
        p.level = depth + item_depth
        run = p.add_run()
        run.text = prefix
        run.font.size = Pt(14)
        _emit_inline_runs(p, text, base_size=Pt(14))
    return True


def _b_quote(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    text = _str(block.get("text"))
    cite = _str(block.get("cite"))
    if not text and not cite:
        return False
    if text:
        for line in text.splitlines() or [""]:
            p = _next_paragraph(frame)
            p.level = depth
            run = p.add_run()
            run.text = f"“{line}”"
            run.font.italic = True
            run.font.size = Pt(14)
    if cite:
        p = _next_paragraph(frame)
        p.level = depth
        run = p.add_run()
        run.text = f"— {cite}"
        run.font.size = Pt(12)
    return True


def _b_callout(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    variant = _str(block.get("variant")) or "info"
    title = _str(block.get("title"))
    text = _str(block.get("text"))
    if not (title or text):
        return False
    p = _next_paragraph(frame)
    p.level = depth
    tag_run = p.add_run()
    tag_run.text = f"[{variant.upper()}] "
    tag_run.font.bold = True
    tag_run.font.size = Pt(14)
    if title:
        title_run = p.add_run()
        title_run.text = title
        title_run.font.bold = True
        title_run.font.size = Pt(14)
    if text:
        for line in text.splitlines():
            sub_p = _next_paragraph(frame)
            sub_p.level = depth + 1
            _emit_inline_runs(sub_p, line, base_size=Pt(13))
    return True


def _b_code(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    code = _str(block.get("code"))
    language = _str(block.get("language"))
    filename = _str(block.get("filename"))
    if filename or language:
        p = _next_paragraph(frame)
        p.level = depth
        run = p.add_run()
        run.text = filename or language or "code"
        run.font.bold = True
        run.font.size = Pt(12)
    if not code:
        return True
    for line in code.splitlines() or [""]:
        p = _next_paragraph(frame)
        p.level = depth
        run = p.add_run()
        run.text = line or " "
        run.font.name = "Consolas"
        run.font.size = Pt(11)
    return True


def _b_math(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    expr = _str(block.get("expression"))
    if not expr:
        return False
    p = _next_paragraph(frame)
    p.level = depth
    run = p.add_run()
    run.text = expr
    run.font.name = "Cambria Math"
    run.font.italic = True
    run.font.size = Pt(14)
    return True


def _b_table(slide: Any, _frame: Any, block: dict[str, Any], _ctx: _Ctx, _depth: int) -> bool:
    headers = block.get("headers") or []
    rows = block.get("rows") or []
    if not headers and not rows:
        return False
    if not headers:
        cols = max((len(r) for r in rows), default=0)
        headers = [""] * cols
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    if n_cols == 0:
        return False
    left = _BODY_LEFT
    top = _next_shape_top(slide)
    width = _BODY_WIDTH
    height = Inches(min(0.4 * n_rows + 0.4, 5.2))
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = _str(h)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(12)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx in range(n_cols):
            cell = table.cell(r_idx, c_idx)
            cell.text = _str(row[c_idx]) if c_idx < len(row) else ""
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
    return True


def _b_kpi_cards(slide: Any, _frame: Any, block: dict[str, Any], _ctx: _Ctx, _depth: int) -> bool:
    items = block.get("items") or []
    if not items:
        return False
    # 2x2 grid (truncate beyond 4).
    grid = items[:4]
    top = _next_shape_top(slide)
    cell_w = Inches(5.8)
    cell_h = Inches(1.8)
    gap = Inches(0.2)
    for i, it in enumerate(grid):
        col = i % 2
        row = i // 2
        left = _BODY_LEFT + (cell_w + gap) * col
        cur_top = top + (cell_h + gap) * row
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, cur_top, cell_w, cell_h)
        box.fill.solid()
        box.fill.fore_color.rgb = _hex_rgb("F1F5F9")
        box.line.color.rgb = _hex_rgb("CBD5E1")
        tf = box.text_frame
        tf.word_wrap = True
        label_p = tf.paragraphs[0]
        label_p.alignment = PP_ALIGN.LEFT
        run = label_p.add_run()
        run.text = _str(it.get("label"))
        run.font.size = Pt(12)
        value_p = tf.add_paragraph()
        run = value_p.add_run()
        run.text = _str(it.get("value"))
        run.font.size = Pt(24)
        run.font.bold = True
        delta = _str(it.get("delta"))
        trend = _str(it.get("trend"))
        if delta or trend:
            d_p = tf.add_paragraph()
            run = d_p.add_run()
            run.text = f"{delta} ({trend})" if delta and trend else delta or trend
            run.font.size = Pt(11)
    return True


def _b_chart(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    chart_type = _str(block.get("chartType") or block.get("chart_type")).lower()
    title = _str(block.get("title"))
    data = block.get("data") or {}
    labels = data.get("labels") or []
    series = data.get("series") or []
    xl_type = {
        "line": XL_CHART_TYPE.LINE,
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "pie": XL_CHART_TYPE.PIE,
    }.get(chart_type)
    if xl_type is None or not labels or not series:
        # Fallback: textual placeholder + summary.
        p = _next_paragraph(frame)
        p.level = depth
        run = p.add_run()
        run.text = f"[차트] {title or chart_type or 'chart'}".strip()
        run.font.bold = True
        run.font.size = Pt(14)
        for sd in series:
            sub_p = _next_paragraph(frame)
            sub_p.level = depth + 1
            run = sub_p.add_run()
            name = _str(sd.get("name"))
            values = sd.get("values") or []
            run.text = f"{name}: {', '.join(_str(v) for v in values)}"
            run.font.size = Pt(12)
        return True
    chart_data = CategoryChartData()
    chart_data.categories = [_str(label) for label in labels]
    for sd in series:
        chart_data.add_series(_str(sd.get("name")) or "", [_to_number(v) for v in (sd.get("values") or [])])
    left = _BODY_LEFT
    top = _next_shape_top(slide)
    width = Inches(11)
    height = Inches(4.5)
    chart_shape = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart = chart_shape.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    return True


def _b_gantt(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    tasks = block.get("tasks") or []
    if not tasks:
        return False
    p = _next_paragraph(frame)
    p.level = depth
    run = p.add_run()
    run.text = "[Gantt]"
    run.font.bold = True
    run.font.size = Pt(14)
    for t in tasks:
        sub_p = _next_paragraph(frame)
        sub_p.level = depth + 1
        run = sub_p.add_run()
        run.text = (
            f"{_str(t.get('name'))} : {_str(t.get('start'))} → {_str(t.get('end'))}"
        )
        run.font.size = Pt(12)
    return True


def _b_flow(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    engine = _str(block.get("engine")) or "mermaid"
    source = _str(block.get("source"))
    p = _next_paragraph(frame)
    p.level = depth
    run = p.add_run()
    run.text = f"[{engine} 다이어그램]"
    run.font.bold = True
    run.font.size = Pt(14)
    if source:
        for line in source.splitlines():
            sub_p = _next_paragraph(frame)
            sub_p.level = depth + 1
            run = sub_p.add_run()
            run.text = line
            run.font.name = "Consolas"
            run.font.size = Pt(11)
    return True


def _b_org_chart(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    root = block.get("root") or {}

    def render_node(node: dict[str, Any], cur_depth: int) -> None:
        label = _str(node.get("label"))
        role = _str(node.get("role"))
        line = f"{label} — {role}" if role else label
        if not line:
            return
        p = _next_paragraph(frame)
        p.level = depth + cur_depth
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        for child in node.get("children") or []:
            render_node(child, cur_depth + 1)

    if not root:
        return False
    p = _next_paragraph(frame)
    p.level = depth
    run = p.add_run()
    run.text = "[조직도]"
    run.font.bold = True
    run.font.size = Pt(14)
    render_node(root, 1)
    return True


def _b_iframe(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    src = _str(block.get("src"))
    title = _str(block.get("title")) or src
    return _emit_link_block(frame, "[iframe] ", title, src, depth)


def _b_video(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    url = _str(block.get("url"))
    title = _str(block.get("title")) or url
    return _emit_link_block(frame, "🎬 ", title, url, depth)


def _b_image(slide: Any, frame: Any, block: dict[str, Any], ctx: _Ctx, depth: int) -> bool:
    image_id = _str(block.get("imageId") or block.get("image_id"))
    caption = _str(block.get("caption"))
    resolver = ctx.opts.image_resolver
    info: dict[str, Any] | None = None
    if resolver and image_id:
        try:
            info = resolver(image_id)
        except Exception:
            info = None
    blob = info.get("bytes") if info else None
    if not blob:
        return _emit_link_block(
            frame, "🖼 ", caption or image_id or "image", f"image:{image_id}", depth
        )
    left = _BODY_LEFT
    top = _next_shape_top(slide)
    width = Inches(8)
    try:
        slide.shapes.add_picture(io.BytesIO(blob), left, top, width=width)
    except Exception:
        return _emit_link_block(
            frame, "🖼 ", caption or image_id or "image", f"image:{image_id}", depth
        )
    if caption:
        p = _next_paragraph(frame)
        p.level = depth
        run = p.add_run()
        run.text = caption
        run.font.italic = True
        run.font.size = Pt(11)
    return True


def _b_gallery(slide: Any, frame: Any, block: dict[str, Any], ctx: _Ctx, depth: int) -> bool:
    items = block.get("items") or []
    if not items:
        return False
    emitted = False
    for it in items:
        single = {
            "type": "image",
            "imageId": it.get("imageId") or it.get("image_id"),
            "caption": it.get("caption"),
            "alt": it.get("alt"),
        }
        if _b_image(slide, frame, single, ctx, depth):
            emitted = True
    return emitted


def _b_file(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    file_id = _str(block.get("fileId") or block.get("file_id"))
    name = _str(block.get("name")) or file_id or "file"
    href = f"/api/v1/files/{file_id}/download" if file_id else ""
    return _emit_link_block(frame, "📎 ", name, href, depth)


def _b_doc_link_card(
    slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int
) -> bool:
    slug = _str(block.get("slug"))
    title = _str(block.get("title")) or slug
    return _emit_link_block(frame, "📄 ", title, f"/docs/{slug}" if slug else "", depth)


def _b_glossary_ref(
    slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int
) -> bool:
    term = _str(block.get("term"))
    if not term:
        return False
    p = _next_paragraph(frame)
    p.level = depth
    run = p.add_run()
    run.text = term
    run.font.italic = True
    run.font.size = Pt(13)
    return True


def _b_columns(slide: Any, frame: Any, block: dict[str, Any], ctx: _Ctx, depth: int) -> bool:
    columns = block.get("columns") or []
    emitted = False
    for col in columns:
        for inner in col:
            if _render_block_into(slide, frame, inner, ctx, depth=depth):
                emitted = True
    return emitted


def _b_tabs(slide: Any, frame: Any, block: dict[str, Any], ctx: _Ctx, depth: int) -> bool:
    tabs = block.get("tabs") or []
    emitted = False
    for tab in tabs:
        label = _str(tab.get("label"))
        if label:
            p = _next_paragraph(frame)
            p.level = depth
            run = p.add_run()
            run.text = f"▸ {label}"
            run.font.bold = True
            run.font.size = Pt(14)
        for inner in tab.get("blocks") or []:
            if _render_block_into(slide, frame, inner, ctx, depth=depth + 1):
                emitted = True
    return emitted


def _b_accordion(slide: Any, frame: Any, block: dict[str, Any], ctx: _Ctx, depth: int) -> bool:
    items = block.get("items") or []
    emitted = False
    for it in items:
        label = _str(it.get("label"))
        if label:
            p = _next_paragraph(frame)
            p.level = depth
            run = p.add_run()
            run.text = f"▾ {label}"
            run.font.bold = True
            run.font.size = Pt(14)
        for inner in it.get("blocks") or []:
            if _render_block_into(slide, frame, inner, ctx, depth=depth + 1):
                emitted = True
    return emitted


def _b_data_source(
    slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int
) -> bool:
    endpoint = _str(block.get("endpoint"))
    render = _str(block.get("render"))
    p = _next_paragraph(frame)
    p.level = depth
    run = p.add_run()
    run.text = f"[데이터 소스] {endpoint} · render={render}"
    run.font.size = Pt(13)
    return True


def _b_dashboard_embed(
    slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int
) -> bool:
    provider = _str(block.get("provider"))
    panel = _str(block.get("panelId") or block.get("panel_id"))
    return _emit_link_block(
        frame, "[대시보드] ", f"{provider} · panel={panel}", "", depth
    )


def _b_calculator(
    slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int
) -> bool:
    label = _str(block.get("label"))
    formula = _str(block.get("formula"))
    if not (label or formula):
        return False
    p = _next_paragraph(frame)
    p.level = depth
    run = p.add_run()
    run.text = f"[계산기] {label}".strip()
    run.font.bold = True
    run.font.size = Pt(14)
    if formula:
        sub_p = _next_paragraph(frame)
        sub_p.level = depth + 1
        run = sub_p.add_run()
        run.text = f"수식: {formula}"
        run.font.name = "Consolas"
        run.font.size = Pt(12)
    return True


def _b_unknown(slide: Any, frame: Any, block: dict[str, Any], _ctx: _Ctx, depth: int) -> bool:
    btype = _str(block.get("type"))
    p = _next_paragraph(frame)
    p.level = depth
    run = p.add_run()
    run.text = f"[알 수 없는 블록: {btype}]"
    run.font.italic = True
    run.font.size = Pt(12)
    return True


_BLOCK_HANDLERS: dict[str, Any] = {
    "paragraph": _b_paragraph,
    "heading-4": _b_heading_4,
    "list": _b_list,
    "quote": _b_quote,
    "callout": _b_callout,
    "code": _b_code,
    "math": _b_math,
    "table": _b_table,
    "kpi-cards": _b_kpi_cards,
    "chart": _b_chart,
    "gantt": _b_gantt,
    "flow": _b_flow,
    "org-chart": _b_org_chart,
    "iframe": _b_iframe,
    "video": _b_video,
    "image": _b_image,
    "gallery": _b_gallery,
    "file": _b_file,
    "doc-link-card": _b_doc_link_card,
    "glossary-ref": _b_glossary_ref,
    "columns": _b_columns,
    "tabs": _b_tabs,
    "accordion": _b_accordion,
    "data-source": _b_data_source,
    "dashboard-embed": _b_dashboard_embed,
    "calculator": _b_calculator,
}


# ── Helpers ──────────────────────────────────────────────────────────


def _str(v: Any) -> str:
    if v is None:
        return ""
    if isinstance(v, bool):
        return "true" if v else "false"
    return str(v)


def _to_number(v: Any) -> float:
    try:
        return float(v)
    except (TypeError, ValueError):
        return 0.0


def _next_paragraph(frame: Any) -> Any:
    """Use the first (empty) paragraph if untouched, else append a new one.

    text_frame always has one paragraph by default. Reusing it for the first
    block keeps the body looking right.
    """
    if (
        len(frame.paragraphs) == 1
        and not frame.paragraphs[0].runs
        and not _str(frame.paragraphs[0].text)
    ):
        return frame.paragraphs[0]
    return frame.add_paragraph()


def _next_shape_top(slide: Any) -> Any:
    """Pick a top offset below already-placed shapes so we don't overlap.

    Title + body box always exist (added in :func:`_new_content_slide` /
    :func:`_body_text_frame`). Walk the rest and place the new shape under
    the lowest one.
    """
    bottom = _BODY_TOP + Inches(0.4)
    for shape in slide.shapes:
        if shape.has_text_frame and shape.top == _TITLE_TOP:
            continue
        try:
            shape_bottom = shape.top + shape.height
        except Exception:
            continue
        if shape_bottom > bottom:
            bottom = shape_bottom + Emu(120000)  # ~0.13in gap
    return bottom


def _emit_inline_runs(p: Any, text: str, *, base_size: Any) -> None:
    """Append runs honouring **bold**, *italic*, ~~strike~~ tokens.

    python-pptx 의 ``Font`` 는 ``strike`` 를 직접 노출하지 않으므로
    underlying ``rPr`` 엘리먼트에 ``strike="sngStrike"`` 속성을 직접 단다.
    """
    pos = 0
    for m in _INLINE_RE.finditer(text):
        if m.start() > pos:
            run = p.add_run()
            run.text = text[pos : m.start()]
            run.font.size = base_size
        token = m.group(0)
        run = p.add_run()
        if token.startswith("**"):
            run.text = m.group(2)
            run.font.bold = True
        elif token.startswith("~~"):
            run.text = m.group(4)
            _set_run_strike(run)
        else:
            run.text = m.group(3)
            run.font.italic = True
        run.font.size = base_size
        pos = m.end()
    if pos < len(text):
        run = p.add_run()
        run.text = text[pos:]
        run.font.size = base_size


def _set_run_strike(run: Any) -> None:
    """``a:rPr/@strike="sngStrike"`` — python-pptx 미노출 속성 직접 패치."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("strike", "sngStrike")


def _emit_link_block(frame: Any, prefix: str, label: str, url: str, depth: int) -> bool:
    if not (label or url):
        return False
    p = _next_paragraph(frame)
    p.level = depth
    run = p.add_run()
    run.text = f"{prefix}{label}"
    run.font.size = Pt(13)
    if url:
        url_p = _next_paragraph(frame)
        url_p.level = depth + 1
        run = url_p.add_run()
        run.text = url
        run.font.size = Pt(11)
        run.font.name = "Consolas"
    return True


def _hex_rgb(hex_str: str) -> Any:
    from pptx.dml.color import RGBColor

    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _append_speaker_note(slide: Any, text: str) -> None:
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    if tf.text:
        tf.text = f"{tf.text}\n{text}"
    else:
        tf.text = text


def _approx_char_count(section: dict[str, Any]) -> int:
    """Cheap character-count of a subsection's textual content.

    Used only for the inline-vs-spawn decision; doesn't have to be exact.
    """
    total = 0

    def count_blocks(blocks: list[dict[str, Any]]) -> None:
        nonlocal total
        for b in blocks or []:
            for k in ("text", "title", "code", "expression", "source"):
                v = b.get(k)
                if isinstance(v, str):
                    total += len(v)
            for it in b.get("items") or []:
                if isinstance(it, str):
                    total += len(it)
                elif isinstance(it, dict):
                    total += len(_str(it.get("text") or it.get("label") or ""))
            for row in b.get("rows") or []:
                if isinstance(row, list):
                    for cell in row:
                        total += len(_str(cell))

    count_blocks(section.get("blocks") or [])
    for sub in section.get("subsections") or []:
        total += _approx_char_count(sub)
    return total
