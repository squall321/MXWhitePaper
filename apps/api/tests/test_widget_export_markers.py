"""Export-side widget marker emission tests.

Verifies that ``emit_marker_text`` returns the correct ``"Widget: <type>"`` text
for widgets that need an export-side marker prepend, returns ``None`` for the
auto-detected opt-out widgets, and that ``render_docx`` actually emits those
markers as plain-text paragraphs for round-trip recognition.
"""
from __future__ import annotations

import io

import ulid
from docx import Document
from pptx import Presentation

from app.services.docx_export import DocxOptions, render_docx
from app.services.pptx_export import PptxOptions, render_pptx
from app.services.widget_markers import emit_marker_text


def _u() -> str:
    return str(ulid.new())


def _build_doc(blocks: list[dict]) -> dict:
    return {
        "schema_version": "1.0",
        "id": _u(),
        "slug": "wm-export-test",
        "title": "T",
        "metadata": {
            "division": "MX",
            "owners": ["t@e.com"],
            "tags": [],
            "confidentiality": "internal",
        },
        "sections": [
            {
                "id": _u(),
                "number": "1",
                "level": 1,
                "title": "sec",
                "blocks": blocks,
                "subsections": [],
            }
        ],
    }


def _all_paragraphs(blob: bytes) -> list[str]:
    return [p.text for p in Document(io.BytesIO(blob)).paragraphs]


def _all_pptx_text(blob: bytes) -> str:
    pres = Presentation(io.BytesIO(blob))
    chunks: list[str] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # python-pptx stub keeps text_frame on the concrete subclass.
                for para in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                    for run in para.runs:
                        chunks.append(run.text)
                    chunks.append(para.text)
    return "\n".join(chunks)


# ── emit_marker_text unit tests ──────────────────────────────────────


def test_emit_marker_for_iframe() -> None:
    assert emit_marker_text({"type": "iframe", "src": "http://x"}) == "Widget: iframe"


def test_emit_marker_for_chart_with_variant() -> None:
    assert (
        emit_marker_text({"type": "chart", "chartType": "line"})
        == "Widget: chart (line)"
    )


def test_emit_marker_for_chart_invalid_chartType_no_variant() -> None:
    assert (
        emit_marker_text({"type": "chart", "chartType": "unknown"})
        == "Widget: chart"
    )


def test_emit_marker_for_doc_link_card_uses_marker_key() -> None:
    assert (
        emit_marker_text({"type": "doc-link-card", "slug": "x"})
        == "Widget: doc-link"
    )


def test_emit_marker_for_glossary_ref_uses_marker_key() -> None:
    assert (
        emit_marker_text({"type": "glossary-ref", "term": "x"})
        == "Widget: glossary"
    )


def test_emit_marker_for_callout_uses_variant() -> None:
    # Round-trip guarantee: every widget emits a marker so the docx
    # importer can reconstruct it even when the visible body (coloured
    # paragraph / 2x2 grid / multi-column / image series) isn't self-
    # identifying. variant carried verbatim when in the allowed set.
    assert (
        emit_marker_text({"type": "callout", "variant": "warn", "text": "x"})
        == "Widget: callout (warn)"
    )


def test_emit_marker_for_callout_invalid_variant_emits_bare() -> None:
    # Unknown variant → no variant token. Importer falls back to info.
    assert (
        emit_marker_text({"type": "callout", "variant": "rainbow", "text": "x"})
        == "Widget: callout"
    )


def test_emit_marker_for_kpi_cards_is_bare() -> None:
    assert emit_marker_text({"type": "kpi-cards"}) == "Widget: kpi-cards"


def test_emit_marker_for_gallery_grid_is_bare() -> None:
    # Default layout "grid" is implied — no variant token.
    assert emit_marker_text({"type": "gallery", "layout": "grid"}) == "Widget: gallery"


def test_emit_marker_for_gallery_carousel_has_variant() -> None:
    assert (
        emit_marker_text({"type": "gallery", "layout": "carousel"})
        == "Widget: gallery (carousel)"
    )


def test_emit_marker_for_columns_encodes_count_as_variant() -> None:
    assert (
        emit_marker_text({"type": "columns", "columns": [[], []]})
        == "Widget: columns (2)"
    )
    assert (
        emit_marker_text({"type": "columns", "columns": [[], [], []]})
        == "Widget: columns (3)"
    )


def test_emit_marker_returns_none_for_paragraph() -> None:
    assert emit_marker_text({"type": "paragraph"}) is None


def test_emit_marker_returns_none_for_table() -> None:
    assert emit_marker_text({"type": "table"}) is None


# ── docx render integration tests ────────────────────────────────────


def test_docx_export_emits_marker_for_iframe() -> None:
    doc = _build_doc(
        [{"type": "iframe", "id": _u(), "src": "https://example.com"}]
    )
    blob = render_docx(doc, options=DocxOptions())
    paragraphs = _all_paragraphs(blob)
    assert "Widget: iframe" in paragraphs


def test_docx_export_emits_marker_for_chart() -> None:
    doc = _build_doc(
        [
            {
                "type": "chart",
                "id": _u(),
                "chartType": "line",
                "data": {
                    "labels": ["a", "b"],
                    "series": [{"name": "s1", "values": [1, 2]}],
                },
            }
        ]
    )
    blob = render_docx(doc, options=DocxOptions())
    paragraphs = _all_paragraphs(blob)
    assert "Widget: chart (line)" in paragraphs


def test_docx_export_emits_hidden_marker_for_callout() -> None:
    """Callout now emits a hidden marker for lossless round-trip. The
    coloured paragraph remains the visible representation; the marker
    rides above it with run.font.hidden=True (invisible in normal Word
    view, picked up by docx_import via parse_marker)."""
    doc = _build_doc(
        [{"type": "callout", "id": _u(), "variant": "warn", "text": "hi"}]
    )
    blob = render_docx(doc, options=DocxOptions())
    document = Document(io.BytesIO(blob))
    found = False
    for p in document.paragraphs:
        for r in p.runs:
            if r.font.hidden and "Widget: callout" in r.text:
                found = True
                break
    assert found, "Hidden callout marker missing"


def test_docx_export_emits_hidden_marker_for_kpi_cards() -> None:
    doc = _build_doc(
        [
            {
                "type": "kpi-cards",
                "id": _u(),
                "items": [{"label": "L", "value": "V"}],
            }
        ]
    )
    blob = render_docx(doc, options=DocxOptions())
    document = Document(io.BytesIO(blob))
    found = False
    for p in document.paragraphs:
        for r in p.runs:
            if r.font.hidden and "Widget: kpi-cards" in r.text:
                found = True
                break
    assert found, "Hidden kpi-cards marker missing"


# ── pptx render integration tests ────────────────────────────────────


def test_pptx_export_emits_marker_for_iframe() -> None:
    doc = _build_doc(
        [{"type": "iframe", "id": _u(), "src": "https://example.com"}]
    )
    blob = render_pptx(doc, options=PptxOptions())
    text = _all_pptx_text(blob)
    assert "Widget: iframe" in text


def test_pptx_export_emits_marker_for_chart_with_variant() -> None:
    doc = _build_doc(
        [
            {
                "type": "chart",
                "id": _u(),
                "chartType": "bar",
                "data": {
                    "labels": ["a", "b"],
                    "series": [{"name": "s1", "values": [1, 2]}],
                },
            }
        ]
    )
    blob = render_pptx(doc, options=PptxOptions())
    text = _all_pptx_text(blob)
    assert "Widget: chart (bar)" in text


def test_pptx_export_emits_marker_for_doc_link_card() -> None:
    doc = _build_doc(
        [{"type": "doc-link-card", "slug": "x", "id": _u()}]
    )
    blob = render_pptx(doc, options=PptxOptions())
    text = _all_pptx_text(blob)
    assert "Widget: doc-link" in text


def test_pptx_export_no_marker_for_callout() -> None:
    doc = _build_doc(
        [{"type": "callout", "id": _u(), "variant": "warn", "text": "hi"}]
    )
    blob = render_pptx(doc, options=PptxOptions())
    text = _all_pptx_text(blob)
    assert "Widget: callout" not in text


# ── html render integration tests ────────────────────────────────────


def test_html_export_emits_comment_marker_for_iframe() -> None:
    from app.services.html_renderer import render_namuwiki_html

    doc = _build_doc(
        [{"type": "iframe", "id": _u(), "src": "https://example.com"}]
    )
    out = render_namuwiki_html(doc)
    assert "<!-- Widget: iframe -->" in out


def test_html_export_emits_comment_for_chart_with_variant() -> None:
    from app.services.html_renderer import render_namuwiki_html

    doc = _build_doc(
        [
            {
                "type": "chart",
                "id": _u(),
                "chartType": "line",
                "data": {
                    "labels": ["a", "b"],
                    "series": [{"name": "s1", "values": [1, 2]}],
                },
            }
        ]
    )
    out = render_namuwiki_html(doc)
    assert "<!-- Widget: chart (line) -->" in out


def test_html_export_no_marker_for_callout() -> None:
    from app.services.html_renderer import render_namuwiki_html

    doc = _build_doc(
        [{"type": "callout", "id": _u(), "variant": "warn", "text": "hi"}]
    )
    out = render_namuwiki_html(doc)
    assert "<!-- Widget: callout" not in out


# ── markdown render integration tests ────────────────────────────────


def test_markdown_export_emits_marker_for_iframe() -> None:
    from app.services.markdown_export import render_markdown

    doc = _build_doc(
        [{"type": "iframe", "id": _u(), "src": "https://example.com"}]
    )
    out = render_markdown(doc)
    assert "Widget: iframe" in out


def test_markdown_export_emits_marker_for_chart_with_variant() -> None:
    from app.services.markdown_export import render_markdown

    doc = _build_doc(
        [
            {
                "type": "chart",
                "id": _u(),
                "chartType": "line",
                "data": {
                    "labels": ["a", "b"],
                    "series": [{"name": "s1", "values": [1, 2]}],
                },
            }
        ]
    )
    out = render_markdown(doc)
    assert "Widget: chart (line)" in out


def test_markdown_export_no_marker_for_callout() -> None:
    from app.services.markdown_export import render_markdown

    doc = _build_doc(
        [{"type": "callout", "id": _u(), "variant": "warn", "text": "hi"}]
    )
    out = render_markdown(doc)
    assert "Widget: callout" not in out


# ── pdf / whiteboard / image-annotation: marker + body across renderers ──


def _pdf_block() -> dict:
    return {
        "type": "pdf",
        "id": _u(),
        "file_id": _u(),
        "title": "Spec v1",
        "page": 3,
    }


def _whiteboard_block() -> dict:
    return {
        "type": "whiteboard",
        "id": _u(),
        "viewbox": {"w": 1000, "h": 600},
        "elements": [],
    }


def _image_annotation_block() -> dict:
    return {
        "type": "image-annotation",
        "id": _u(),
        "image_id": _u(),
        "annotations": [],
    }


def test_pptx_export_emits_marker_and_body_for_pdf() -> None:
    doc = _build_doc([_pdf_block()])
    blob = render_pptx(doc, options=PptxOptions())
    text = _all_pptx_text(blob)
    assert "Widget: pdf" in text
    assert "Spec v1" in text
    assert "📕" in text


def test_pptx_export_emits_marker_and_body_for_whiteboard() -> None:
    doc = _build_doc([_whiteboard_block()])
    blob = render_pptx(doc, options=PptxOptions())
    text = _all_pptx_text(blob)
    assert "Widget: whiteboard" in text
    assert "🖼" in text
    assert "1000" in text and "600" in text


def test_pptx_export_emits_marker_and_body_for_image_annotation() -> None:
    block = _image_annotation_block()
    doc = _build_doc([block])
    blob = render_pptx(doc, options=PptxOptions())
    text = _all_pptx_text(blob)
    assert "Widget: image-annotation" in text
    assert "🖼" in text
    assert block["image_id"] in text


def test_html_export_emits_marker_and_body_for_pdf() -> None:
    from app.services.html_renderer import render_namuwiki_html

    doc = _build_doc([_pdf_block()])
    out = render_namuwiki_html(doc)
    assert "<!-- Widget: pdf -->" in out
    assert "widget-pdf" in out
    assert "Spec v1" in out


def test_html_export_emits_marker_and_body_for_whiteboard() -> None:
    from app.services.html_renderer import render_namuwiki_html

    doc = _build_doc([_whiteboard_block()])
    out = render_namuwiki_html(doc)
    assert "<!-- Widget: whiteboard -->" in out
    assert "widget-whiteboard" in out
    assert "1000" in out and "600" in out


def test_html_export_emits_marker_and_body_for_image_annotation() -> None:
    from app.services.html_renderer import render_namuwiki_html

    block = _image_annotation_block()
    doc = _build_doc([block])
    out = render_namuwiki_html(doc)
    assert "<!-- Widget: image-annotation -->" in out
    assert "widget-image-annotation" in out
    assert block["image_id"] in out


def test_markdown_export_emits_marker_and_body_for_pdf() -> None:
    from app.services.markdown_export import render_markdown

    doc = _build_doc([_pdf_block()])
    out = render_markdown(doc)
    assert "Widget: pdf" in out
    assert "PDF" in out
    assert "Spec v1" in out


def test_markdown_export_emits_marker_and_body_for_whiteboard() -> None:
    from app.services.markdown_export import render_markdown

    doc = _build_doc([_whiteboard_block()])
    out = render_markdown(doc)
    assert "Widget: whiteboard" in out
    assert "Whiteboard" in out
    assert "1000" in out and "600" in out


def test_markdown_export_emits_marker_and_body_for_image_annotation() -> None:
    from app.services.markdown_export import render_markdown

    block = _image_annotation_block()
    doc = _build_doc([block])
    out = render_markdown(doc)
    assert "Widget: image-annotation" in out
    assert "Annotated image" in out
    assert block["image_id"] in out


def test_b_gallery_emits_hidden_marker() -> None:
    blocks = [{
        "type": "gallery", "id": _u(), "layout": "grid",
        "items": [{"imageId": _u()}, {"imageId": _u()}, {"imageId": _u()}],
    }]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    found = False
    for p in doc.paragraphs:
        for r in p.runs:
            if r.font.hidden and "Widget: gallery" in r.text:
                found = True
    assert found


def test_b_gallery_carousel_variant_in_marker() -> None:
    blocks = [{
        "type": "gallery", "id": _u(), "layout": "carousel",
        "items": [{"imageId": _u()}, {"imageId": _u()}, {"imageId": _u()}],
    }]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    found_variant = False
    for p in doc.paragraphs:
        for r in p.runs:
            if r.font.hidden and "Widget: gallery (carousel)" in r.text:
                found_variant = True
    assert found_variant


# ── callout hidden-marker round-trip ─────────────────────────────────


def _walk_all(sections: list[dict]) -> list[dict]:
    out: list[dict] = []
    for s in sections or []:
        out.extend(s.get("blocks") or [])
        out.extend(_walk_all(s.get("subsections") or []))
    return out


def test_b_callout_emits_hidden_marker() -> None:
    """A CalloutBlock rendered to docx should carry a hidden 'Widget: callout (variant)'
    marker as the first paragraph, so docx → docx_to_document round-trip
    reconstructs the CalloutBlock instead of a plain paragraph."""
    blocks = [
        {"type": "callout", "id": _u(), "variant": "warn", "text": "주의"}
    ]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    # The first paragraph SHOULD be the hidden marker (or near the top).
    found_marker = False
    for p in doc.paragraphs:
        for r in p.runs:
            if r.font.hidden and "Widget: callout" in r.text:
                found_marker = True
                break
        if found_marker:
            break
    assert found_marker, "Hidden callout marker missing from docx"


def test_callout_roundtrips_via_hidden_marker() -> None:
    """docx_export → docx_to_document should reconstruct CalloutBlock via the
    hidden marker — proves the round-trip loop closes for callouts."""
    from app.services.docx_import import docx_to_document

    blocks = [
        {"type": "callout", "id": _u(), "variant": "warn", "text": "주의"}
    ]
    blob = render_docx(_build_doc(blocks))
    result = docx_to_document(blob, slug="rt", title="", owner_user_id=_u())
    out = _walk_all(result["document"]["sections"])
    callouts = [b for b in out if b.get("type") == "callout"]
    assert callouts, [b.get("type") for b in out]
    assert callouts[0]["variant"] == "warn"
    assert "주의" in callouts[0]["text"]


# ── kpi-cards hidden-marker round-trip ───────────────────────────────


def test_b_kpi_cards_emits_hidden_marker() -> None:
    blocks = [{
        "type": "kpi-cards", "id": _u(),
        "items": [{"label": "매출", "value": "1억"}, {"label": "MAU", "value": "5만"}],
    }]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    found = False
    for p in doc.paragraphs:
        for r in p.runs:
            if r.font.hidden and "Widget: kpi-cards" in r.text:
                found = True
                break
    assert found


def test_kpi_cards_roundtrips_via_hidden_marker() -> None:
    from app.services.docx_import import docx_to_document

    blocks = [{
        "type": "kpi-cards", "id": _u(),
        "items": [{"label": "매출", "value": "1억"}, {"label": "MAU", "value": "5만"}],
    }]
    blob = render_docx(_build_doc(blocks))
    result = docx_to_document(blob, slug="rt", title="", owner_user_id=_u())
    out = _walk_all(result["document"]["sections"])
    kpi_blocks = [b for b in out if b.get("type") == "kpi-cards"]
    assert kpi_blocks, [b.get("type") for b in out]
    labels = [it["label"] for it in kpi_blocks[0]["items"]]
    assert "매출" in labels
    assert "MAU" in labels


def test_b_kpi_cards_preserves_delta_and_trend() -> None:
    from app.services.docx_import import docx_to_document

    blocks = [{
        "type": "kpi-cards", "id": _u(),
        "items": [{"label": "L", "value": "V", "delta": "+10%", "trend": "up"}],
    }]
    blob = render_docx(_build_doc(blocks))
    result = docx_to_document(blob, slug="rt", title="", owner_user_id=_u())
    out = _walk_all(result["document"]["sections"])
    kpi = next(b for b in out if b.get("type") == "kpi-cards")
    assert kpi["items"][0].get("delta") == "+10%"
    assert kpi["items"][0].get("trend") == "up"


# ── image-annotation / whiteboard hidden-marker checks ───────────────


def test_b_image_annotation_marker_is_hidden() -> None:
    blocks = [{
        "type": "image-annotation", "id": _u(),
        "image_id": _u(), "annotations": [],
    }]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    for p in doc.paragraphs:
        for r in p.runs:
            if "Widget: image-annotation" in r.text:
                assert r.font.hidden is True
                return
    raise AssertionError("Marker not found at all")


def test_b_whiteboard_marker_is_hidden() -> None:
    blocks = [{
        "type": "whiteboard", "id": _u(),
        "viewbox": {"w": 800, "h": 600}, "elements": [],
    }]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    for p in doc.paragraphs:
        for r in p.runs:
            if "Widget: whiteboard" in r.text:
                assert r.font.hidden is True
                return
    raise AssertionError("Marker not found at all")


# ── Sanity guard: every marker-emitting widget must use hidden text ──


def _minimal_block(kind: str) -> dict | None:
    """Return a minimal-valid DocumentJSON block for ``kind`` whose export
    triggers ``emit_marker_text``. Returns None for types without an easy
    fixture (none currently — every marker-emitting type is covered)."""
    if kind == "iframe":
        return {"type": "iframe", "id": _u(), "src": "https://example.com"}
    if kind == "video":
        return {"type": "video", "id": _u(),
                "url": "https://www.youtube.com/watch?v=abc",
                "provider": "youtube"}
    if kind == "file":
        return {"type": "file", "id": _u(), "fileId": _u(), "name": "x.docx"}
    if kind == "pdf":
        return {"type": "pdf", "id": _u(), "file_id": _u(), "title": "Spec"}
    if kind == "whiteboard":
        return {"type": "whiteboard", "id": _u(),
                "viewbox": {"w": 800, "h": 600}, "elements": []}
    if kind == "image-annotation":
        return {"type": "image-annotation", "id": _u(),
                "image_id": _u(), "annotations": []}
    if kind == "flow":
        return {"type": "flow", "id": _u(),
                "engine": "mermaid", "source": "graph TD; A-->B"}
    if kind == "chart":
        return {"type": "chart", "id": _u(), "chartType": "bar",
                "data": {"labels": ["a"],
                         "series": [{"name": "s1", "values": [1]}]}}
    if kind == "gantt":
        return {"type": "gantt", "id": _u(),
                "tasks": [{"name": "T", "start": "2026-01-01",
                           "end": "2026-01-02"}]}
    if kind == "org-chart":
        return {"type": "org-chart", "id": _u(),
                "root": {"id": _u(), "label": "Root"}}
    if kind == "tabs":
        return {"type": "tabs", "id": _u(),
                "tabs": [{"label": "T1", "blocks": []}]}
    if kind == "accordion":
        return {"type": "accordion", "id": _u(),
                "items": [{"label": "I1", "blocks": []}]}
    if kind == "doc-link-card":
        return {"type": "doc-link-card", "id": _u(), "slug": "x"}
    if kind == "glossary-ref":
        return {"type": "glossary-ref", "id": _u(), "term": "Hyperloop"}
    # kpi-cards / callout / gallery / columns aren't in _EXPORT_MARKER_TYPES,
    # but if some other agent adds them later this stays forward-compatible:
    if kind == "kpi-cards":
        return {"type": "kpi-cards", "id": _u(),
                "items": [{"label": "L", "value": "V"}]}
    if kind == "callout":
        return {"type": "callout", "id": _u(),
                "variant": "warn", "text": "hi"}
    if kind == "gallery":
        return {"type": "gallery", "id": _u(),
                "items": [{"imageId": _u()}]}
    if kind == "columns":
        return {"type": "columns", "id": _u(), "columns": [
            [{"type": "paragraph", "id": _u(), "text": "a"}],
            [{"type": "paragraph", "id": _u(), "text": "b"}],
        ]}
    return None


def test_all_marker_emitting_widgets_use_hidden_text() -> None:
    """Sanity: every widget that emits a marker via emit_marker_text MUST
    do so as hidden text. Regression guard against accidentally reverting
    to visible markers."""
    from app.services.widget_markers import _EXPORT_MARKER_TYPES
    expected_hidden = sorted(_EXPORT_MARKER_TYPES)

    # Build minimal-valid blocks for each marker-emitting type and render.
    for kind in expected_hidden:
        block = _minimal_block(kind)
        if block is None:
            continue  # type with no easy minimal fixture
        blob = render_docx(_build_doc([block]))
        doc = Document(io.BytesIO(blob))
        marker_runs = []
        for p in doc.paragraphs:
            for r in p.runs:
                if r.text.startswith("Widget:"):
                    marker_runs.append((kind, r.font.hidden, r.text))
        for kind_t, hidden, text in marker_runs:
            assert hidden is True, f"{kind_t}: marker '{text}' is VISIBLE, must be hidden"


# ── columns hidden-marker round-trip ─────────────────────────────────


def test_b_columns_emits_hidden_marker() -> None:
    blocks = [{
        "type": "columns", "id": _u(),
        "columns": [
            [{"type": "paragraph", "id": _u(), "text": "left col"}],
            [{"type": "paragraph", "id": _u(), "text": "right col"}],
        ],
    }]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    found = False
    for p in doc.paragraphs:
        for r in p.runs:
            if r.font.hidden and "Widget: columns" in r.text:
                found = True
    assert found


def test_columns_roundtrips_via_hidden_marker() -> None:
    from app.services.docx_import import docx_to_document

    blocks = [{
        "type": "columns", "id": _u(),
        "columns": [
            [{"type": "paragraph", "id": _u(), "text": "왼쪽 단"}],
            [{"type": "paragraph", "id": _u(), "text": "오른쪽 단"}],
        ],
    }]
    blob = render_docx(_build_doc(blocks))
    result = docx_to_document(blob, slug="rt", title="", owner_user_id=_u())
    out = _walk_all(result["document"]["sections"])
    cols = [b for b in out if b.get("type") == "columns"]
    assert cols, [b.get("type") for b in out]
    assert len(cols[0]["columns"]) >= 2


# ── pass-2 MED gaps: pdf page / org-chart layout / glossary-ref ──────


def test_b_pdf_page_hidden_marker_emitted_when_non_default() -> None:
    """pass-2 M3 — `Widget: pdf` marker doesn't carry the page number, so
    docx_export emits a separate hidden `⟦pdf:page=N⟧` run when page != 1
    (the default) so round-trip can later recover the page reference."""
    blocks = [{
        "type": "pdf",
        "id": _u(),
        "file_id": _u(),
        "title": "Sample doc",
        "page": 5,
    }]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    found_pdf_marker = False
    found_page_marker = False
    for p in doc.paragraphs:
        for r in p.runs:
            if r.font.hidden and "Widget: pdf" in r.text:
                found_pdf_marker = True
            if r.font.hidden and r.text == "⟦pdf:page=5⟧":
                found_page_marker = True
    assert found_pdf_marker
    assert found_page_marker


def test_b_pdf_page_hidden_marker_skipped_when_default() -> None:
    """page=1 is the default — no `⟦pdf:page=…⟧` hidden marker should be
    emitted. The standard `Widget: pdf` marker is still present."""
    blocks = [{
        "type": "pdf",
        "id": _u(),
        "file_id": _u(),
        "title": "Sample doc",
        "page": 1,
    }]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    for p in doc.paragraphs:
        for r in p.runs:
            assert not (r.font.hidden and r.text.startswith("⟦pdf:page="))


def test_b_org_chart_horizontal_layout_hidden_marker() -> None:
    """pass-2 M6 — org-chart `Widget: org-chart` marker doesn't carry a
    layout variant in the standard grammar, so docx_export emits a separate
    hidden `⟦org-chart:layout=horizontal⟧` run for non-default layouts."""
    blocks = [{
        "type": "org-chart",
        "id": _u(),
        "layout": "horizontal",
        "root": {
            "id": "n1",
            "label": "CEO",
            "children": [
                {"id": "n2", "label": "CTO", "children": []},
                {"id": "n3", "label": "CFO", "children": []},
            ],
        },
    }]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    found_layout = False
    for p in doc.paragraphs:
        for r in p.runs:
            if r.font.hidden and r.text == "⟦org-chart:layout=horizontal⟧":
                found_layout = True
    assert found_layout


def test_b_glossary_ref_renders_without_definition_field() -> None:
    """pass-2 M11 regression — GlossaryRefBlock schema only has `term`; the
    legacy docx_export code that tried to read `block.get("definition")` is
    dead and removed. Export must still produce a term-only paragraph."""
    blocks = [{
        "type": "glossary-ref",
        "id": _u(),
        "term": "ULID",
    }]
    blob = render_docx(_build_doc(blocks))
    doc = Document(io.BytesIO(blob))
    has_term = False
    for p in doc.paragraphs:
        if "ULID" in p.text and not all(r.font.hidden for r in p.runs):
            has_term = True
    assert has_term
