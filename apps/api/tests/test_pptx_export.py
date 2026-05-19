"""PPTX export 렌더러 단위 테스트.

renderer 자체는 부수효과 없는 pure function. endpoint 통합은 한 건만 스모크.
.pptx 는 zip 컨테이너이므로:
  - magic bytes (PK\\x03\\x04) 로 zip 여부 검증
  - python-pptx 로 다시 열어 슬라이드 수 / 본문 텍스트 확인
"""
from __future__ import annotations

import io
import zipfile

import pytest
from httpx import ASGITransport, AsyncClient
from pptx import Presentation
from pptx.presentation import Presentation as PptxPresentation

from app.main import app
from app.services.pptx_export import PptxOptions, render_pptx


def _doc(blocks: list[dict] | None = None, **overrides: object) -> dict:
    base = {
        "schema_version": "1.0",
        "id": "01TESTDOC0000000000000000Z",
        "slug": "fixture-pptx",
        "title": "발표용 픽스처",
        "summary": "단위 테스트 픽스처입니다.",
        "metadata": {
            "division": "MX",
            "team": "Editor",
            "owners": ["someone@example.com"],
            "tags": ["unit", "pptx"],
            "confidentiality": "internal",
        },
        "sections": [
            {
                "id": "01SEC00000000000000000000A",
                "number": "1",
                "level": 1,
                "title": "서론",
                "blocks": blocks or [
                    {
                        "type": "paragraph",
                        "id": "01P000000000000000000000A1",
                        "text": "첫 문단입니다.",
                    }
                ],
                "subsections": [],
            }
        ],
    }
    base.update(overrides)
    return base


def _slides_text(prs: PptxPresentation) -> str:
    """모든 슬라이드의 모든 text frame 을 직렬화 (검색용)."""
    buf: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            # python-pptx stub keeps `text_frame` / `table` on the concrete
            # subclasses only; `has_text_frame`/`has_table` are runtime guards.
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                    for run in para.runs:
                        buf.append(run.text)
            if shape.has_table:
                table = shape.table  # type: ignore[attr-defined]
                for row in table.rows:
                    for cell in row.cells:
                        buf.append(cell.text)
    return "\n".join(buf)


# ── pure renderer ────────────────────────────────────────────────────


def test_renderer_returns_valid_pptx_zip() -> None:
    out = render_pptx(_doc())
    # PK\x03\x04 — zip local file header.
    assert out[:4] == b"PK\x03\x04"
    # zipfile must accept it as a valid archive.
    with zipfile.ZipFile(io.BytesIO(out)) as zf:
        names = zf.namelist()
        # Standard pptx parts.
        assert any(n.endswith("/presentation.xml") for n in names)


def test_renderer_emits_title_slide_plus_section_slides() -> None:
    doc = _doc()
    doc["sections"].append(
        {
            "id": "01SEC00000000000000000001A",
            "number": "2",
            "level": 1,
            "title": "본론",
            "blocks": [
                {
                    "type": "paragraph",
                    "id": "01P00000000000000000001A2",
                    "text": "두 번째 섹션입니다.",
                }
            ],
            "subsections": [],
        }
    )
    out = render_pptx(doc)
    prs = Presentation(io.BytesIO(out))
    # title + 2 sections = 3 slides
    assert len(prs.slides) == 3
    text = _slides_text(prs)
    assert "발표용 픽스처" in text
    assert "1 서론" in text
    assert "2 본론" in text
    assert "두 번째 섹션입니다." in text


def test_renderer_paragraph_inline_bold_italic_strike() -> None:
    blocks = [
        {
            "type": "paragraph",
            "id": "01P000000000000000000010A",
            "text": "이것은 **굵게** 와 *기울임* 그리고 ~~취소~~ 입니다.",
        }
    ]
    out = render_pptx(_doc(blocks))
    prs = Presentation(io.BytesIO(out))
    # Walk runs and check formatting was applied to the right text fragments.
    found_bold = False
    found_italic = False
    found_strike = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                for run in para.runs:
                    if run.text == "굵게" and run.font.bold:
                        found_bold = True
                    if run.text == "기울임" and run.font.italic:
                        found_italic = True
                    # python-pptx exposes strike via XML; .font.strike is the API.
                    if run.text == "취소":
                        # strike attr may be None when unset; rendered runs
                        # should set it explicitly to True.
                        if getattr(run.font, "_rPr", None) is not None or True:
                            # fall back to checking the XML element directly
                            xml = run._r.xml
                            if 'strike="sngStrike"' in xml or "strike=" in xml:
                                found_strike = True
    assert found_bold
    assert found_italic
    assert found_strike


def test_renderer_table_emits_native_table_shape() -> None:
    blocks = [
        {
            "type": "table",
            "id": "01T1A0000000000000000001A",
            "headers": ["이름", "수량"],
            "rows": [["사과", "3"], ["배", "2"]],
        }
    ]
    out = render_pptx(_doc(blocks))
    prs = Presentation(io.BytesIO(out))
    table_shapes = [
        s for slide in prs.slides for s in slide.shapes if s.has_table
    ]
    assert len(table_shapes) == 1
    table = table_shapes[0].table  # type: ignore[attr-defined]
    assert table.cell(0, 0).text == "이름"
    assert table.cell(0, 1).text == "수량"
    assert table.cell(1, 0).text == "사과"
    assert table.cell(2, 1).text == "2"


def test_renderer_chart_line_emits_native_chart() -> None:
    blocks = [
        {
            "type": "chart",
            "id": "01CH0000000000000000001A",
            "chartType": "line",
            "title": "분기 매출",
            "data": {
                "labels": ["Q1", "Q2", "Q3"],
                "series": [{"name": "rev", "values": [1, 2, 3]}],
            },
        }
    ]
    out = render_pptx(_doc(blocks))
    prs = Presentation(io.BytesIO(out))
    chart_shapes = [
        s for slide in prs.slides for s in slide.shapes if s.has_chart
    ]
    assert len(chart_shapes) == 1


def test_renderer_chart_unsupported_type_falls_back_to_text() -> None:
    blocks = [
        {
            "type": "chart",
            "id": "01CH0000000000000000002A",
            "chartType": "scatter",
            "title": "산점도",
            "data": {
                "labels": ["A", "B"],
                "series": [{"name": "s", "values": [1, 2]}],
            },
        }
    ]
    out = render_pptx(_doc(blocks))
    prs = Presentation(io.BytesIO(out))
    text = _slides_text(prs)
    assert "[차트]" in text
    assert "산점도" in text


def test_renderer_kpi_cards_emit_2x2_grid() -> None:
    blocks = [
        {
            "type": "kpi-cards",
            "id": "01K10000000000000000001A",
            "items": [
                {"label": "매출", "value": "100", "delta": "+5", "trend": "up"},
                {"label": "이익", "value": "30", "delta": "-1", "trend": "down"},
                {"label": "고객", "value": "1.2k"},
                {"label": "NPS", "value": "42"},
            ],
        }
    ]
    out = render_pptx(_doc(blocks))
    prs = Presentation(io.BytesIO(out))
    text = _slides_text(prs)
    assert "매출" in text and "100" in text
    assert "이익" in text and "30" in text
    assert "고객" in text and "NPS" in text


def test_renderer_speaker_note_appended_to_slide_notes() -> None:
    blocks = [
        {
            "type": "paragraph",
            "id": "01P00000000000000000020A",
            "text": "본문",
            "meta": {"note": "speaker: 여기서 천천히 강조합니다."},
        }
    ]
    out = render_pptx(_doc(blocks))
    prs = Presentation(io.BytesIO(out))
    # Slide 0 = title, slide 1 = our section.
    notes_tf = prs.slides[1].notes_slide.notes_text_frame
    assert notes_tf is not None  # speaker note added above
    note = notes_tf.text
    assert "여기서 천천히 강조합니다." in note


def test_renderer_long_subsection_spawns_separate_slide() -> None:
    doc = _doc()
    long_text = "긴 본문 " * 200  # > 500 chars
    doc["sections"][0]["subsections"] = [
        {
            "id": "01SUB000000000000000000001",
            "number": "1.1",
            "level": 2,
            "title": "롱 서브섹션",
            "blocks": [
                {
                    "type": "paragraph",
                    "id": "01P00000000000000000030A",
                    "text": long_text,
                }
            ],
            "subsections": [],
        }
    ]
    out = render_pptx(doc)
    prs = Presentation(io.BytesIO(out))
    # title + section + spawned sub = 3 slides
    assert len(prs.slides) == 3
    text = _slides_text(prs)
    assert "1.1 롱 서브섹션" in text


def test_renderer_short_subsection_inlines_as_bullets() -> None:
    doc = _doc()
    doc["sections"][0]["subsections"] = [
        {
            "id": "01SUB000000000000000000002",
            "number": "1.1",
            "level": 2,
            "title": "짧은 서브",
            "blocks": [
                {
                    "type": "paragraph",
                    "id": "01P00000000000000000040A",
                    "text": "짧은 본문",
                }
            ],
            "subsections": [],
        }
    ]
    out = render_pptx(doc)
    prs = Presentation(io.BytesIO(out))
    # title + section only — subsection inlined.
    assert len(prs.slides) == 2
    text = _slides_text(prs)
    assert "1.1 짧은 서브" in text
    assert "짧은 본문" in text


def test_renderer_image_uses_resolver_bytes_when_available() -> None:
    # 1x1 PNG (smallest valid image bytes).
    png = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cc\xfa\xcf"
        b"\x00\x00\x00\x02\x00\x01\xe2!\xbc3\x00\x00\x00\x00IEND\xaeB`\x82"
    )

    def resolver(image_id: str) -> dict | None:
        if image_id == "01IMG00000000000000000000":
            return {"bytes": png, "mime": "image/png"}
        return None

    blocks = [
        {
            "type": "image",
            "id": "01I10000000000000000001A",
            "imageId": "01IMG00000000000000000000",
            "caption": "그림 1",
        }
    ]
    out = render_pptx(_doc(blocks), options=PptxOptions(image_resolver=resolver))
    prs = Presentation(io.BytesIO(out))
    pics = [
        s
        for slide in prs.slides
        for s in slide.shapes
        if s.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pics) == 1


def test_renderer_image_without_resolver_falls_back_to_text() -> None:
    blocks = [
        {
            "type": "image",
            "id": "01I10000000000000000002A",
            "imageId": "01IMG00000000000000000999",
            "caption": "그림 2",
        }
    ]
    out = render_pptx(_doc(blocks))
    prs = Presentation(io.BytesIO(out))
    text = _slides_text(prs)
    assert "그림 2" in text


def test_renderer_list_styles_bullet_number_check() -> None:
    blocks = [
        {
            "type": "list",
            "id": "01L1A0000000000000000001A",
            "style": "number",
            "items": ["하나", "둘"],
        },
        {
            "type": "list",
            "id": "01L1B0000000000000000001A",
            "style": "check",
            "items": ["TODO1", "TODO2"],
        },
    ]
    out = render_pptx(_doc(blocks))
    prs = Presentation(io.BytesIO(out))
    text = _slides_text(prs)
    assert "1. " in text and "2. " in text
    assert "하나" in text and "둘" in text
    assert "☐" in text
    assert "TODO1" in text


def test_renderer_bibliography_block_emits_title_and_entries() -> None:
    """BibliographyBlock → heading paragraph + one paragraph per entry."""
    blocks = [
        {
            "type": "bibliography",
            "id": "01BIB00000000000000000001",
            "title": "참고",
            "entries": [
                {"key": "smith2020", "text": "Smith, J. (2020). Foo.", "url": "https://example.org/foo"},
                {"text": "익명 보고서, 2021."},
            ],
        }
    ]
    out = render_pptx(_doc(blocks))
    prs = Presentation(io.BytesIO(out))
    text = _slides_text(prs)
    assert "참고" in text
    assert "Smith, J. (2020). Foo." in text
    assert "[smith2020]" in text
    assert "익명 보고서, 2021." in text
    assert "https://example.org/foo" in text


# ── endpoint integration ─────────────────────────────────────────────


SEED_SLUG = "month-end-closing"


@pytest.mark.asyncio
async def test_export_pptx_endpoint_returns_attachment() -> None:
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        r = await ac.post("/api/v1/exports/pptx", json={"slug": SEED_SLUG})
    assert r.status_code == 200, r.text
    ctype = r.headers.get("content-type", "")
    assert "presentationml.presentation" in ctype
    cd = r.headers.get("content-disposition") or ""
    assert "attachment" in cd
    assert f"{SEED_SLUG}.pptx" in cd
    body = r.content
    assert body[:4] == b"PK\x03\x04"


@pytest.mark.asyncio
async def test_export_pptx_endpoint_404_for_missing_slug() -> None:
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        r = await ac.post(
            "/api/v1/exports/pptx", json={"slug": "no-such-slug-xxx-xxx"}
        )
    assert r.status_code == 404


@pytest.mark.asyncio
async def test_export_pptx_endpoint_rejects_missing_slug() -> None:
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        r = await ac.post("/api/v1/exports/pptx", json={})
    assert r.status_code == 422
